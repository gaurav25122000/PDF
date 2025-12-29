from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import Response
from fastapi.middleware.cors import CORSMiddleware
from mangum import Mangum
from pypdf import PdfReader, PdfWriter
from io import BytesIO
from PIL import Image
from reportlab.lib.utils import ImageReader

app = FastAPI()

# CORS configuration
origins = [
    "http://localhost:5173",
    "http://localhost:8888",
    "*" # Adjust for production
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
async def root():
    return {"message": "iLovePDF Clone API is running"}

@app.get("/api/health")
async def health_check():
    return {"status": "ok"}

from pypdf import PdfWriter, PdfReader
from io import BytesIO
from fastapi import UploadFile, HTTPException
from fastapi.responses import Response

@app.post("/api/process/merge")
async def merge_pdfs(files: list[UploadFile]):
    if not files or len(files) < 2:
        raise HTTPException(status_code=400, detail="At least 2 PDF files are required for merging.")
    
    try:
        writer = PdfWriter()
        
        for file in files:
            # Read file into memory
            content = await file.read()
            pdf_content = BytesIO(content)
            reader = PdfReader(pdf_content)
            
            # Append pages
            for page in reader.pages:
                writer.add_page(page)
        
        # Write merged PDF to memory
        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="merged_document.pdf"'
        }
        
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
        
    except Exception as e:
        print(f"Error merging PDFs: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error merging PDFs: {str(e)}")

@app.post("/api/process/split")
async def split_pdf(file: UploadFile, pages: str):
    """
    Split/Extract pages from a PDF.
    pages format: "1,2,5-7" (1-based indexing)
    Returns: A single PDF containing the selected pages.
    """
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()
        
        total_pages = len(reader.pages)
        selected_pages = set()
        
        # Parse page ranges
        parts = pages.split(',')
        for part in parts:
            part = part.strip()
            if '-' in part:
                start, end = map(int, part.split('-'))
                # 1-based to 0-based, inclusive
                for p in range(start - 1, end):
                    if 0 <= p < total_pages:
                        selected_pages.add(p)
            else:
                p = int(part) - 1
                if 0 <= p < total_pages:
                    selected_pages.add(p)
        
        sorted_pages = sorted(list(selected_pages))
        
        if not sorted_pages:
             raise HTTPException(status_code=400, detail="No valid pages selected.")

        for p in sorted_pages:
            writer.add_page(reader.pages[p])
            
        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="split_document.pdf"'
        }
        
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
        
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid page range format.")
    except Exception as e:
        print(f"Error splitting PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error splitting PDF: {str(e)}")

@app.post("/api/process/protect")
async def protect_pdf(file: UploadFile, password: str):
    if not file or not password:
        raise HTTPException(status_code=400, detail="PDF file and password are required.")
    
    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()
        
        writer.append_pages_from_reader(reader)
        writer.encrypt(password)
        
        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="protected_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error protecting PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error protecting PDF: {str(e)}")

@app.post("/api/process/unlock")
async def unlock_pdf(file: UploadFile, password: str):
    if not file or not password:
        raise HTTPException(status_code=400, detail="PDF file and password are required.")
    
    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        
        if reader.is_encrypted:
            try:
                reader.decrypt(password)
            except:
                raise HTTPException(status_code=400, detail="Incorrect password.")
        else:
             # If not encrypted, just return original? Or error?
             # iLovePDF Unlock usually removes security. If not secured, it's a no-op.
             pass

        writer = PdfWriter()
        writer.append_pages_from_reader(reader)
        
        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="unlocked_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error unlocking PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error unlocking PDF: {str(e)}")

@app.post("/api/process/compress")
async def compress_pdf(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()
        
        writer.append_pages_from_reader(reader)
        # Attempt to compress
        try:
             # Basic pypdf compression (removes duplicates, etc)
             writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)
             
             # Also can loop pages and compress content streams
             for page in writer.pages:
                 page.compress_content_streams()
        except:
             pass 

        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="compressed_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error compressing PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error compressing PDF: {str(e)}")

from PIL import Image

@app.post("/api/process/jpg-to-pdf")
async def jpg_to_pdf(files: list[UploadFile]):
    if not files:
        raise HTTPException(status_code=400, detail="At least one image file is required.")

    try:
        image_list = []
        for file in files:
            content = await file.read()
            img = Image.open(BytesIO(content)).convert("RGB")
            image_list.append(img)
        
        if not image_list:
            raise HTTPException(status_code=400, detail="No valid images found.")

        output_pdf = BytesIO()
        # Save first image and append the rest
        image_list[0].save(output_pdf, save_all=True, append_images=image_list[1:], format="PDF")
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="converted_images.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error converting JPG to PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting JPG to PDF: {str(e)}")

@app.post("/api/process/rotate")
async def rotate_pdf(file: UploadFile, angle: int = 90):
    """
    Rotate all pages by the specified angle (must be multiple of 90: 90, 180, 270).
    """
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    if angle % 90 != 0:
        raise HTTPException(status_code=400, detail="Angle must be a multiple of 90.")

    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)
            
        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="rotated_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error rotating PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error rotating PDF: {str(e)}")

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

@app.post("/api/process/watermark")
async def watermark_pdf(file: UploadFile, text: str = "Watermark"):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()

        # Create watermark PDF in memory
        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        
        # Draw transparent text
        can.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.5)
        can.setFont("Helvetica", 50)
        # Center approx
        can.saveState()
        can.translate(300, 400)
        can.rotate(45)
        can.drawCentredString(0, 0, text)
        can.restoreState()
        can.save()
        
        packet.seek(0)
        watermark_pdf = PdfReader(packet)
        watermark_page = watermark_pdf.pages[0]

        for page in reader.pages:
            # Merge watermark page into current page
            page.merge_page(watermark_page)
            writer.add_page(page)
            
        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="watermarked_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error adding watermark: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error adding watermark: {str(e)}")

@app.post("/api/process/page-numbers")
async def page_numbers_pdf(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()
        
        total_pages = len(reader.pages)

        for i, page in enumerate(reader.pages):
            packet = BytesIO()
            # Use page size from original if possible, else default
            # Getting page size from pypdf page obj is a bit tricky for reportlab canvas init
            # Simplified: Assuming standard letter or using huge canvas and placing relative
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            
            can = canvas.Canvas(packet, pagesize=(width, height))
            can.setFont("Helvetica", 10)
            text = f"Page {i+1} of {total_pages}"
            # Draw at bottom center
            can.drawCentredString(width / 2, 20, text)
            can.save()
            
            packet.seek(0)
            number_pdf = PdfReader(packet)
            page.merge_page(number_pdf.pages[0])
            writer.add_page(page)

        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="numbered_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
    except Exception as e:
        print(f"Error adding page numbers: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error adding page numbers: {str(e)}")

import fitz # PyMuPDF
import zipfile
import os
import tempfile
from pdf2docx import Converter

@app.post("/api/process/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        content = await file.read()
        # Open PDF from bytes
        doc = fitz.open(stream=content, filetype="pdf")
        
        output_zip = BytesIO()
        with zipfile.ZipFile(output_zip, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                # Render page to image (pixmap)
                pix = page.get_pixmap(dpi=150) # 150 DPI is decent balance
                img_data = pix.tobytes("jpg")
                zf.writestr(f"page_{i+1}.jpg", img_data)
        
        output_zip.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="converted_images.zip"'
        }
        return Response(content=output_zip.getvalue(), media_type="application/zip", headers=headers)
    except Exception as e:
        print(f"Error converting PDF to JPG: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting PDF to JPG: {str(e)}")

@app.post("/api/process/pdf-to-word")
async def pdf_to_word(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    tmp_pdf = None
    tmp_docx = None
    try:
        # pdf2docx requires file paths nicely
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f:
            f.write(await file.read())
            tmp_pdf = f.name
            
        tmp_docx = tmp_pdf.replace(".pdf", ".docx")
        
        cv = Converter(tmp_pdf)
        cv.convert(tmp_docx)
        cv.close()
        
        with open(tmp_docx, "rb") as f:
            docx_content = f.read()
            
        headers = {
            'Content-Disposition': 'attachment; filename="converted_document.docx"'
        }
        return Response(content=docx_content, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers=headers)
    except Exception as e:
        print(f"Error converting PDF to Word: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting PDF to Word: {str(e)}")
    finally:
        # Cleanup temp files
        if tmp_pdf and os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
        if tmp_docx and os.path.exists(tmp_docx):
            os.remove(tmp_docx)

import json
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
import base64

@app.post("/api/process/edit")
async def edit_pdf(file: UploadFile, operations: str = Form(...)):
    """
    Apply edits (text, overlay images) to a PDF.
    operations: JSON string of list of objects:
    [
        {
            "page": 0,
            "type": "text",
            "text": "Hello",
            "x": 100,
            "y": 100, # Bottom-left origin (PDF standard)
            "fontSize": 12,
            "color": "#000000"
        },
        {
            "page": 0,
            "type": "image",
            "data": "base64...",
            "x": 50,
            "y": 50,
            "width": 100,
            "height": 100
        }
    ]
    """
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        ops_list = json.loads(operations)
    except:
        raise HTTPException(status_code=400, detail="Invalid operations JSON.")

    try:
        content = await file.read()
        pdf_content = BytesIO(content)
        reader = PdfReader(pdf_content)
        writer = PdfWriter()
        
        # We need to process page by page
        # It's more efficient to group ops by page
        ops_by_page = {}
        for op in ops_list:
            p_idx = op.get("page", 0)
            if p_idx not in ops_by_page:
                ops_by_page[p_idx] = []
            ops_by_page[p_idx].append(op)
            
        for i, page in enumerate(reader.pages):
            if i in ops_by_page:
                # Create detailed overlay for this page
                packet = BytesIO()
                # Get page size to ensure overlay matches
                width = float(page.mediabox.width)
                height = float(page.mediabox.height)
                
                c = canvas.Canvas(packet, pagesize=(width, height))
                
                for op in ops_by_page[i]:
                    x = float(op.get("x", 0))
                    # Frontend usually gives Top-Left origin coordinates. PDF uses Bottom-Left.
                    # We might need to invert Y if frontend sends TL coords.
                    # Let's assume frontend sends PDF-ready BL coords OR we handle it here.
                    # DECISION: We interpret 'x' and 'y' as Top-Left coordinates from the frontend (0,0 is top-left).
                    # PDF uses Bottom-Left origin (0,0 is bottom-left).
                    # So conversion: new_y = page_height - y_from_top
                    
                    y_top = float(op.get("y", 0))
                    y = height - y_top
                    
                    if op.get("type") == "text":
                        text = op.get("text", "")
                        f_size = int(op.get("fontSize", 12))
                        color_hex = op.get("color", "#000000")
                        c.setFont("Helvetica", f_size)
                        c.setFillColor(HexColor(color_hex))
                        c.drawString(x, y, text)
                        
                    elif op.get("type") == "image":
                        img_data = op.get("data", "")
                        if img_data:
                            # Handle base64
                            if "," in img_data:
                                img_data = img_data.split(",")[1]
                            img_bytes = base64.b64decode(img_data)
                            img_io = BytesIO(img_bytes)
                            img_pil = Image.open(img_io)
                            
                            w = float(op.get("width", 100))
                            h = float(op.get("height", 100))
                            
                            # drawImage in reportlab
                            # y is bottom-left of image
                            c.drawImage(ImageReader(img_pil), x, y, width=w, height=h, mask='auto')

                c.save()
                packet.seek(0)
                overlay_pdf = PdfReader(packet)
                page.merge_page(overlay_pdf.pages[0])
            
            writer.add_page(page)

        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="edited_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
        
    except Exception as e:
        print(f"Error editing PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error editing PDF: {str(e)}")

    except Exception as e:
        print(f"Error editing PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error editing PDF: {str(e)}")

import pdfplumber
import pandas as pd

@app.post("/api/process/pdf-to-excel")
async def pdf_to_excel(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    tmp_pdf = None
    tmp_xlsx = None
    try:
        # pdfplumber works best with paths
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f:
            f.write(await file.read())
            tmp_pdf = f.name
            
        tmp_xlsx = tmp_pdf.replace(".pdf", ".xlsx")
        
        # Extract tables
        with pdfplumber.open(tmp_pdf) as pdf:
            with pd.ExcelWriter(tmp_xlsx, engine='openpyxl') as writer:
                has_tables = False
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        if table:
                            df = pd.DataFrame(table)
                            # Cleanup: first row as header if it looks like one?
                            # For simplicity, dump raw
                            sheet_name = f"Page{i+1}_Table{j+1}"
                            # Excel sheet names limited to 31 chars
                            sheet_name = sheet_name[:31] 
                            df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
                            has_tables = True
                
                if not has_tables:
                    # Create empty sheet
                    pd.DataFrame(["No tables found"]).to_excel(writer, sheet_name="Info", index=False, header=False)

        with open(tmp_xlsx, "rb") as f:
            xlsx_content = f.read()
            
        headers = {
            'Content-Disposition': 'attachment; filename="converted_tables.xlsx"'
        }
        return Response(content=xlsx_content, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers=headers)

    except Exception as e:
        print(f"Error converting PDF to Excel: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting PDF to Excel: {str(e)}")
    finally:
         if tmp_pdf and os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
         if tmp_xlsx and os.path.exists(tmp_xlsx):
            os.remove(tmp_xlsx)

from pptx import Presentation
from pptx.util import Inches
import mammoth
from xhtml2pdf import pisa

@app.post("/api/process/pdf-to-pptx")
async def pdf_to_pptx(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="PDF file is required.")
    
    try:
        content = await file.read()
        doc = fitz.open(stream=content, filetype="pdf")
        
        prs = Presentation()
        # Clean default slide
        if len(prs.slides) > 0:
            # removing slides is tricky in python-pptx, usually starts empty or default template
            # accessing xml to delete is hard.
            # actually prs = Presentation() usually has 0 slides if no template.
            pass

        # Use blank layout
        blank_slide_layout = prs.slide_layouts[6] 
        
        for page in doc:
            # Render page to image
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            img_io = BytesIO(img_data)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide, covering whole slide?
            # PPTX default size is 10x7.5 inches. 
            # We can adjust constraints or just center.
            # Let's just add image fit to slide
            slide.shapes.add_picture(img_io, Left=0, Top=0, Width=prs.slide_width, Height=prs.slide_height)

        output_pptx = BytesIO()
        prs.save(output_pptx)
        output_pptx.seek(0)
        
        headers = {
            'Content-Disposition': 'attachment; filename="converted_presentation.pptx"'
        }
        return Response(content=output_pptx.getvalue(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)

    except Exception as e:
        print(f"Error converting PDF to PPTX: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting PDF to PPTX: {str(e)}")

@app.post("/api/process/word-to-pdf")
async def word_to_pdf(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="Word file is required.")
    
    tmp_docx = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as f:
            f.write(await file.read())
            tmp_docx = f.name
            
        # Convert DOCX -> HTML (using mammoth) -> PDF (using xhtml2pdf)
        with open(tmp_docx, "rb") as doctor:
            result = mammoth.convert_to_html(doctor)
            html = result.value
            
        # Add basic style for better PDF
        html_content = f"""
        <html>
        <head>
        <style>
            body {{ font-family: Helvetica, sans-serif; font-size: 12pt; }}
        </style>
        </head>
        <body>
        {html}
        </body>
        </html>
        """
        
        output_pdf = BytesIO()
        pisa_status = pisa.CreatePDF(html_content, dest=output_pdf)
        
        if pisa_status.err:
             raise Exception("PDF generation failed")
             
        output_pdf.seek(0)

        headers = {
            'Content-Disposition': 'attachment; filename="converted_document.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
        
    except Exception as e:
        print(f"Error converting Word to PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting Word to PDF: {str(e)}")
    finally:
        if tmp_docx and os.path.exists(tmp_docx):
            os.remove(tmp_docx)

@app.post("/api/process/excel-to-pdf")
async def excel_to_pdf(file: UploadFile):
    if not file:
        raise HTTPException(status_code=400, detail="Excel file is required.")
    
    tmp_xlsx = None
    try:
        # Load Excel with Pandas
        content = await file.read()
        tmp_xlsx = BytesIO(content)
        
        # Read all sheets
        # sheet_name=None reads all as dict of dfs
        dfs = pd.read_excel(tmp_xlsx, sheet_name=None)
        
        # Build HTML
        html_parts = []
        html_parts.append("""
        <html>
        <head>
        <style>
            body { font-family: Helvetica, sans-serif; font-size: 10pt; }
            table { border-collapse: collapse; width: 100%; margin-bottom: 20px; }
            th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
            th { background-color: #f2f2f2; }
            h2 { color: #333; }
        </style>
        </head>
        <body>
        """)
        
        for sheet_name, df in dfs.items():
            html_parts.append(f"<h2>{sheet_name}</h2>")
            # Convert DF to HTML
            html_parts.append(df.to_html(index=False, na_rep=""))
            
        html_parts.append("</body></html>")
        
        full_html = "\n".join(html_parts)
        
        output_pdf = BytesIO()
        pisa_status = pisa.CreatePDF(full_html, dest=output_pdf)
        
        if pisa_status.err:
             raise Exception("PDF generation failed")
             
        output_pdf.seek(0)

        headers = {
            'Content-Disposition': 'attachment; filename="converted_spreadsheet.pdf"'
        }
        return Response(content=output_pdf.getvalue(), media_type="application/pdf", headers=headers)
        
    except Exception as e:
        print(f"Error converting Excel to PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error converting Excel to PDF: {str(e)}")

# Serverless handler
handler = Mangum(app)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
